#!/usr/bin/env python
import subprocess, os, sys
from pathlib import Path

def extract_pdf(filepath):
    try:
        import pdfplumber
        with pdfplumber.open(filepath) as pdf:
            pages = [p.extract_text() for p in pdf.pages]
            return chr(10).join(t for t in pages if t)
    except: pass
    try:
        r = subprocess.run(['pdftotext','-layout',filepath,'-'],capture_output=True,text=True,timeout=120)
        return r.stdout
    except: return ''

def extract_video(filepath):
    import whisper
    audio = filepath + '.wav'
    try:
        subprocess.run(['ffmpeg','-y','-i',filepath,'-ar','16000','-ac','1','-vn',audio],capture_output=True,timeout=300)
        m = whisper.load_model('medium')
        return m.transcribe(audio, language='zh')['text']
    finally:
        if os.path.exists(audio): os.remove(audio)

def extract_docx(filepath):
    try:
        from docx import Document
        return chr(10).join(p.text for p in Document(filepath).paragraphs if p.text.strip())
    except: return ''

def extract_pptx(filepath):
    try:
        from pptx import Presentation
        texts = []
        for slide in Presentation(filepath).slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.extend(p.text for p in shape.text_frame.paragraphs if p.text.strip())
        return chr(10).join(texts)
    except: return ''

def extract_txt(filepath):
    for enc in ['utf-8','gbk','gb2312']:
        try: return Path(filepath).read_text(encoding=enc)
        except: pass
    return ''

def extract_xlsx(filepath):
    try:
        import pandas as pd
        df = pd.read_excel(filepath)
        return df.to_csv(index=False)
    except: return ''

def extract_text(filepath):
    ext = Path(filepath).suffix.lower()
    if ext == '.pdf': return extract_pdf(filepath)
    if ext in ('.mp4','.mkv','.avi','.mov','.flv'): return extract_video(filepath)
    if ext in ('.docx','.doc'): return extract_docx(filepath)
    if ext in ('.pptx','.ppt'): return extract_pptx(filepath)
    if ext in ('.txt','.md'): return extract_txt(filepath)
    if ext in ('.xlsx','.xls'): return extract_xlsx(filepath)
    return ''
